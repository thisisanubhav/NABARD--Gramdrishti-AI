"""Builds the GramDrishti AI pitch deck as a real .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ASSETS = os.path.dirname(os.path.abspath(__file__))

# Brand colors
GREEN = RGBColor(0x1E, 0x7A, 0x4C)
CREAM = RGBColor(0xF5, 0xF2, 0xEA)
ORANGE = RGBColor(0xFE, 0x93, 0x2C)
CHARCOAL = RGBColor(0x1C, 0x23, 0x21)
MUTED = RGBColor(0x5C, 0x6B, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_TEXT = RGBColor(0xE8, 0xF2, 0xEC)
BORDER = RGBColor(0xDC, 0xD6, 0xC8)
ROW_ALT = RGBColor(0xEC, 0xE8, 0xDC)

SERIF = "Georgia"
SANS = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_bg(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)
    return rect


def add_text(slide, left, top, width, height, text, size, color, bold=False,
             font=SANS, align=PP_ALIGN.LEFT, italic=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size, color,
                 font=SANS, bold_first=False, space_after=10, bullet_char="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_kicker(slide, text, color):
    add_text(slide, Inches(0.7), Inches(0.4), Inches(6), Inches(0.4),
              text.upper(), 12, color, bold=True, font=SANS)


def add_pic_framed(slide, path, left, top, width, height, border_color=GREEN):
    pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
    pic.line.color.rgb = border_color
    pic.line.width = Pt(1.5)
    return pic


def add_ps_table(slide, left, top, width, header_height, row_height, rows):
    """Manually-drawn Problem/Solution table (avoids fighting pptx's default table styling)."""
    col1_w = Inches(4.3)
    col2_w = width - col1_w
    y = top

    def _cell(text, cx, cw, size, bold, color, valign_h):
        tb = slide.shapes.add_textbox(cx + Inches(0.18), y, cw - Inches(0.32), valign_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = SANS
        run.font.color.rgb = color

    # header row
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, header_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = GREEN
    rect.line.fill.background(); rect.shadow.inherit = False
    _cell(rows[0][0], left, col1_w, 13, True, WHITE, header_height)
    _cell(rows[0][1], left + col1_w, col2_w, 13, True, WHITE, header_height)
    y += header_height

    for i, (prob, sol) in enumerate(rows[1:]):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, row_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = WHITE if i % 2 == 0 else ROW_ALT
        rect.line.color.rgb = BORDER; rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        _cell(prob, left, col1_w, 14, False, CHARCOAL, row_height)
        _cell(sol, left + col1_w, col2_w, 14, True, GREEN, row_height)
        y += row_height


# ---------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN)
add_text(s, Inches(0.9), Inches(0.7), Inches(6), Inches(0.5),
          "NABARD HACKATHON", 14, ORANGE, bold=True, font=SANS)
add_text(s, Inches(0.9), Inches(2.5), Inches(9.5), Inches(1.3),
          "GramDrishti AI", 60, WHITE, bold=True, font=SERIF)
add_text(s, Inches(0.9), Inches(3.55), Inches(8.5), Inches(0.7),
          "AI-powered foresight for rural enterprises", 22, CREAM_TEXT, font=SANS, italic=True)
add_pic_framed(s, os.path.join(ASSETS, "01-login.png"), Inches(7.6), Inches(4.3), Inches(5.2), Inches(2.75))

# Team credential block
divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.75), Inches(1.3), Pt(2.5))
divider.fill.solid(); divider.fill.fore_color.rgb = ORANGE
divider.line.fill.background(); divider.shadow.inherit = False
add_text(s, Inches(0.9), Inches(5.0), Inches(5.5), Inches(0.35),
          "TEAM", 13, ORANGE, bold=True, font=SANS)
add_text(s, Inches(0.9), Inches(5.32), Inches(5.5), Inches(0.6),
          "LoneWolf", 26, WHITE, bold=True, font=SERIF)
add_text(s, Inches(0.9), Inches(5.95), Inches(5.5), Inches(0.45),
          "Anubhav Harsh Sinha", 17, CREAM_TEXT, font=SANS)

add_text(s, Inches(0.9), Inches(6.85), Inches(6), Inches(0.4),
          "Predicting cash-flow risk before it becomes a default", 15, CREAM_TEXT, font=SANS)

# ---------------------------------------------------------------
# SLIDE 2 — Problem -> Solution Map
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "The Problem  →  The Solution", MUTED)
add_text(s, Inches(0.7), Inches(0.72), Inches(11.8), Inches(0.7),
          "Four real problems. Four specific answers.",
          28, CHARCOAL, bold=True, font=SERIF)
add_ps_table(s, Inches(0.7), Inches(1.65), Inches(11.93), Inches(0.5), Inches(1.28), [
    ("PROBLEM", "SOLUTION"),
    ("Risk scores are a black box — no idea why, or what to do",
     "SHAP-explained drivers + a concrete, sector-specific recommendation"),
    ("A field officer can't manually track hundreds of enterprises",
     "Portfolio dashboard: risk distribution, district/sector/climate breakdown, prioritized list"),
    ("Risk is discovered only after a loan default",
     "6-month AI forecast + risk score, recalculated live on every transaction"),
    ("Built for cities, not villages — no offline support, no Hindi, no climate awareness",
     "Offline-first, bilingual, real rainfall/flood/drought data feeding the risk model"),
])

# ---------------------------------------------------------------
# SLIDE 3 — Explainable AI
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Solving Problem 1  ·  The Black Box", ORANGE)
add_text(s, Inches(0.7), Inches(0.85), Inches(6.0), Inches(1.3),
          "A risk score you can act on", 30, CHARCOAL, bold=True, font=SERIF)
add_bullets(s, Inches(0.7), Inches(2.3), Inches(5.8), Inches(4.5), [
    "LightGBM forecasting + SHAP-based explainability — every score comes with its real drivers",
    "“Missed loan repayments — 45%, Loan default on record — 34%, Expenses as share of income — 21%”",
    "Paired with a concrete recommendation (“Talk to your lender early”) — not a diagnosis without a next step",
], 17, CHARCOAL, space_after=20)
add_pic_framed(s, os.path.join(ASSETS, "03-enterprise-detail.png"), Inches(6.9), Inches(0.85), Inches(6.0), Inches(6.0))

# ---------------------------------------------------------------
# SLIDE 4 — Architecture
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Architecture & How It Works", MUTED)
# Source image is 1536x1024 (3:2) — fit to height and center horizontally to avoid distortion
_arch_h = Inches(6.4)
_arch_w = Inches(6.4 * 1536 / 1024)
_arch_left = Emu(int((SW - _arch_w) / 2))
add_pic_framed(s, os.path.join(ASSETS, "09_Architecture_diagram.png"), _arch_left, Inches(0.85), _arch_w, _arch_h)

# ---------------------------------------------------------------
# SLIDE 5 — Field Officer View
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Solving Problem 2  ·  Field Officer View", ORANGE)
add_text(s, Inches(0.7), Inches(0.8), Inches(11.8), Inches(0.8),
          "Manage risk across the entire portfolio", 28, CHARCOAL, bold=True, font=SERIF)
# Source 1600x1100 (1.4545 ratio) — fit to height, keep undistorted
_off_h = Inches(5.5)
_off_w = Inches(5.5 * 1600 / 1100)
add_pic_framed(s, os.path.join(ASSETS, "02-officer-dashboard.png"), Inches(0.55), Inches(1.7), _off_w, _off_h)
_off_text_left = Inches(0.55) + _off_w + Inches(0.45)
_off_text_width = SW - _off_text_left - Inches(0.55)
add_bullets(s, _off_text_left, Inches(2.1), _off_text_width, Inches(4.5), [
    "See all 60 enterprises at a glance — total count and a live high/medium/low risk breakdown",
    "Risk by District and Average Risk by Sector — know exactly where to prioritize a visit",
    "Prioritized intervention list, filterable by sector and risk level, with the AI-driven main factor for each",
    "Real climate risk flags (flood/drought) surfaced automatically for affected enterprises",
], 15, CHARCOAL, space_after=18)

# ---------------------------------------------------------------
# SLIDE 6 — Enterprise Owner View
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Solving Problem 3  ·  Enterprise Owner View", ORANGE)
add_text(s, Inches(0.7), Inches(0.8), Inches(11.8), Inches(0.8),
          "A personal financial co-pilot", 28, CHARCOAL, bold=True, font=SERIF)
# Source 1600x1300 (1.2308 ratio) — fit to height, keep undistorted
_own_h = Inches(5.5)
_own_w = Inches(5.5 * 1600 / 1300)
add_pic_framed(s, os.path.join(ASSETS, "04-owner-dashboard.png"), Inches(0.55), Inches(1.7), _own_w, _own_h)
_own_text_left = Inches(0.55) + _own_w + Inches(0.45)
_own_text_width = SW - _own_text_left - Inches(0.55)
add_bullets(s, _own_text_left, Inches(2.1), _own_text_width, Inches(4.5), [
    "Real KPI cards with month-over-month trend deltas — cash balance, income, expenses, savings",
    "6-month cash-flow forecast with a confidence band, plus a live financial risk score",
    "Risk history timeline shows whether things are improving or worsening over time",
    "Recommendations and transaction entry, right on the same screen",
], 15, CHARCOAL, space_after=18)

# ---------------------------------------------------------------
# SLIDE 6 — What-If Simulator
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Extending Problem 3  ·  From Forecast to Plan", ORANGE)
add_text(s, Inches(0.7), Inches(0.85), Inches(5.8), Inches(1.3),
          "Planning, not just reporting", 30, CHARCOAL, bold=True, font=SERIF)
add_bullets(s, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.0), [
    "Entrepreneurs drag income/expense sliders and watch the 6-month forecast and risk level update live",
    "Turns a static report into a self-service planning tool — “what if my income drops 20%?” answered in real time",
], 18, CHARCOAL, space_after=20)
add_pic_framed(s, os.path.join(ASSETS, "05-owner-simulator-active.png"), Inches(6.6), Inches(0.85), Inches(6.2), Inches(6.0))

# ---------------------------------------------------------------
# SLIDE 7 — Built for Real Conditions & Close
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN)
add_kicker(s, "Solving Problem 4  ·  Built for Real Rural Conditions", ORANGE)
add_bullets(s, Inches(0.7), Inches(1.3), Inches(11.5), Inches(2.6), [
    "Offline-first: transactions queue with no signal and auto-sync when reconnected",
    "Fully bilingual — English and Hindi",
    "Real rainfall/flood/drought data tied to sector and state, feeding directly into risk scoring",
], 19, CREAM_TEXT, space_after=16)
add_text(s, Inches(0.7), Inches(4.3), Inches(11.8), Inches(1.1),
          "Earlier warning → fewer defaults → a stronger rural credit ecosystem.",
          26, WHITE, bold=True, font=SERIF)
add_text(s, Inches(0.7), Inches(5.5), Inches(11.8), Inches(0.6),
          "GramDrishti AI — AI-powered foresight for rural enterprises.", 18, ORANGE, italic=True, font=SANS)

# ---------------------------------------------------------------
# SLIDE 8 — Thank You
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, GREEN)
add_text(s, Inches(0.9), Inches(0.6), Inches(6), Inches(0.5),
          "NABARD HACKATHON", 14, ORANGE, bold=True, font=SANS)
add_text(s, Inches(0.9), Inches(1.3), Inches(11), Inches(1.5),
          "Thank You", 64, WHITE, bold=True, font=SERIF)
add_text(s, Inches(0.9), Inches(2.55), Inches(9), Inches(0.6),
          "GramDrishti AI — AI-powered foresight for rural enterprises", 20, CREAM_TEXT, font=SANS, italic=True)

# Stat recap row (echoes the title slide) — real counts from the demo dataset, not placeholder marketing copy
stats = [("Enterprises Monitored", "60"), ("Sectors Covered", "5"), ("Forecast Horizon", "6 Months")]
for i, (label, value) in enumerate(stats):
    x = Inches(0.9) + Inches(3.3) * i
    add_text(s, x, Inches(3.65), Inches(3.0), Inches(0.35), label.upper(), 12, ORANGE, bold=True, font=SANS)
    add_text(s, x, Inches(4.0), Inches(3.0), Inches(0.8), value, 36, WHITE, bold=True, font=SERIF)

# Team credential block (bookends the title slide)
divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.15), Inches(1.3), Pt(2.5))
divider.fill.solid(); divider.fill.fore_color.rgb = ORANGE
divider.line.fill.background(); divider.shadow.inherit = False
add_text(s, Inches(0.9), Inches(5.4), Inches(5.5), Inches(0.35), "TEAM", 13, ORANGE, bold=True, font=SANS)
add_text(s, Inches(0.9), Inches(5.72), Inches(5.5), Inches(0.6), "LoneWolf", 26, WHITE, bold=True, font=SERIF)
add_text(s, Inches(0.9), Inches(6.35), Inches(6), Inches(0.45), "Anubhav Harsh Sinha", 17, CREAM_TEXT, font=SANS)

add_text(s, Inches(7.2), Inches(6.35), Inches(5.4), Inches(0.6),
          "Submitted for the NABARD Hackathon", 18, ORANGE, italic=True, font=SANS, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------
# SLIDE 9 — Appendix: Technical Deep Dive
# ---------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_kicker(s, "Appendix  ·  Technical Deep Dive", MUTED)
# Source image is 3840x2160 (16:9) — fit to width and center, matches slide ratio closely
_tech_w = Inches(11.2)
_tech_h = Inches(11.2 / (3840 / 2160))
_tech_left = Emu(int((SW - _tech_w) / 2))
add_pic_framed(s, os.path.join(ASSETS, "08-technical-flow-diagram.png"), _tech_left, Inches(0.8), _tech_w, _tech_h)

out_path = os.path.join(ASSETS, "GramDrishti-AI-Pitch-Deck.pptx")
prs.save(out_path)
print("Saved:", out_path)
